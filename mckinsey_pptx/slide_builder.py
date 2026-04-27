"""
McKinsey-style slide builder using python-pptx.
Creates individual slides with proper formatting, action titles, and embedded charts.
"""
from __future__ import annotations
import io
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .config import COLORS, FONTS, LAYOUT, CONTENT_RULES
from .models import SlideContent, SlideType, BulletPoint
from .chart_builder import render_chart
from .native_charts import can_render_native, add_native_chart

# Global flag: set to True to prefer native editable charts over image charts
USE_NATIVE_CHARTS = True


# =============================================================================
# HELPERS
# =============================================================================

def _set_text(text_frame, text: str, font_name: str = FONTS.BODY_FONT,
              font_size=FONTS.BODY_SIZE, bold: bool = False,
              color=COLORS.DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Set text in a text frame with consistent formatting."""
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


def _add_paragraph(text_frame, text: str, font_name: str = FONTS.BODY_FONT,
                   font_size=FONTS.BODY_SIZE, bold: bool = False,
                   color=COLORS.DARK_GRAY, space_after=FONTS.BULLET_SPACING,
                   alignment=PP_ALIGN.LEFT, level: int = 0):
    """Add a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.level = level
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def _add_action_title(slide, title_text: str):
    """Add a McKinsey-style action title (governing thought) to a slide.

    Note: No accent line under the title — accent lines are a known
    hallmark of AI-generated slides. Real McKinsey decks use whitespace
    and the title's bold navy color for visual separation.
    """
    txBox = slide.shapes.add_textbox(
        LAYOUT.TITLE_LEFT, LAYOUT.TITLE_TOP,
        LAYOUT.TITLE_WIDTH, LAYOUT.TITLE_HEIGHT
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, title_text, font_name=FONTS.HEADING_FONT,
              font_size=FONTS.TITLE_SIZE, bold=FONTS.TITLE_BOLD,
              color=FONTS.TITLE_COLOR)


def _add_source_line(slide, source_text: str):
    """Add source citation at the bottom of a slide."""
    txBox = slide.shapes.add_textbox(
        LAYOUT.SOURCE_LEFT, LAYOUT.SOURCE_TOP,
        LAYOUT.SOURCE_WIDTH, LAYOUT.SOURCE_HEIGHT
    )
    tf = txBox.text_frame
    _set_text(tf, f"Source: {source_text}", font_size=FONTS.SMALL_SIZE,
              color=COLORS.LIGHT_GRAY)


def _add_page_number(slide, page_num: int):
    """Add page number at bottom right."""
    txBox = slide.shapes.add_textbox(
        LAYOUT.PAGE_NUM_LEFT, LAYOUT.PAGE_NUM_TOP,
        LAYOUT.PAGE_NUM_WIDTH, Inches(0.3)
    )
    tf = txBox.text_frame
    _set_text(tf, str(page_num), font_size=FONTS.SMALL_SIZE,
              color=COLORS.LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


def _add_bullets(text_frame, bullets: list[BulletPoint]):
    """Add formatted bullet points to a text frame."""
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.space_after = FONTS.BULLET_SPACING
        p.level = 0

        # Bold prefix
        if bullet.bold_prefix:
            run = p.add_run()
            run.text = bullet.bold_prefix + " "
            run.font.name = FONTS.BODY_FONT
            run.font.size = FONTS.BULLET_SIZE
            run.font.bold = True
            run.font.color.rgb = COLORS.DARK_BLUE

        run = p.add_run()
        run.text = bullet.text
        run.font.name = FONTS.BODY_FONT
        run.font.size = FONTS.BULLET_SIZE
        run.font.color.rgb = COLORS.DARK_GRAY

        # Sub-bullets
        for sub in bullet.sub_bullets:
            sp = text_frame.add_paragraph()
            sp.level = 1
            sp.space_after = FONTS.SUB_BULLET_SPACING
            sr = sp.add_run()
            sr.text = sub
            sr.font.name = FONTS.BODY_FONT
            sr.font.size = FONTS.SUB_BULLET_SIZE
            sr.font.color.rgb = FONTS.SUB_BULLET_COLOR


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_title_slide(prs: Presentation, content: SlideContent,
                      client: str = "", date: str = "", confidential: bool = True):
    """Build the title/cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        LAYOUT.WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS.MCKINSEY_BLUE
    bar.line.fill.background()

    # Title — use config size (28pt calibrated from real PDFs, not hardcoded 32pt)
    txBox = slide.shapes.add_textbox(
        LAYOUT.LEFT_MARGIN * 2, Inches(2.2), Inches(10), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, content.action_title, font_name=FONTS.HEADING_FONT,
              font_size=FONTS.TITLE_SIZE, bold=FONTS.TITLE_BOLD, color=FONTS.TITLE_COLOR)

    # Subtitle
    if content.subtitle:
        txBox2 = slide.shapes.add_textbox(
            LAYOUT.LEFT_MARGIN * 2, Inches(3.8), Inches(10), Inches(0.8)
        )
        _set_text(txBox2.text_frame, content.subtitle,
                  font_size=FONTS.SUBTITLE_SIZE, color=COLORS.MEDIUM_GRAY)

    # Client & date
    meta_parts = []
    if client:
        meta_parts.append(client)
    if date:
        meta_parts.append(date)
    if meta_parts:
        txBox3 = slide.shapes.add_textbox(
            LAYOUT.LEFT_MARGIN * 2, Inches(5.0), Inches(10), Inches(0.5)
        )
        _set_text(txBox3.text_frame, " | ".join(meta_parts),
                  font_size=Pt(14), color=COLORS.MEDIUM_GRAY)

    # Confidential notice
    if confidential:
        txBox4 = slide.shapes.add_textbox(
            Inches(1.5), Inches(6.3), Inches(10), Inches(0.4)
        )
        _set_text(txBox4.text_frame, "CONFIDENTIAL AND PROPRIETARY",
                  font_size=Pt(9), bold=True, color=COLORS.LIGHT_GRAY,
                  alignment=PP_ALIGN.LEFT)

    # Bottom bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35),
        LAYOUT.WIDTH, Inches(0.15)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLORS.MCKINSEY_BLUE
    bar2.line.fill.background()

    return slide


def build_executive_summary(prs: Presentation, content: SlideContent, page_num: int = 2):
    """Build an executive summary slide with SCR structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    # Three-column SCR layout
    col_width = Inches(3.7)
    col_gap = Inches(0.4)
    start_left = LAYOUT.CONTENT_LEFT
    top = LAYOUT.CONTENT_TOP

    sections = [
        ("SITUATION", content.situation_text or "", COLORS.MCKINSEY_BLUE),
        ("COMPLICATION", content.complication_text or "", COLORS.NEGATIVE_RED),
        ("RESOLUTION", content.resolution_text or "", COLORS.POSITIVE_GREEN),
    ]

    for i, (label, text, accent_color) in enumerate(sections):
        left = start_left + i * (col_width + col_gap)

        # Section header with colored accent
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, col_width, Pt(3)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = accent_color
        header_bar.line.fill.background()

        # Section label
        txBox = slide.shapes.add_textbox(left, top + Pt(8), col_width, Inches(0.4))
        _set_text(txBox.text_frame, label, font_size=Pt(11), bold=True, color=accent_color)

        # Section content
        txBox2 = slide.shapes.add_textbox(left, top + Inches(0.5), col_width, Inches(4.2))
        txBox2.text_frame.word_wrap = True
        _set_text(txBox2.text_frame, text, font_size=FONTS.EXEC_SUMMARY_SIZE,
                  color=COLORS.DARK_GRAY)

    if content.source:
        _add_source_line(slide, content.source)
    _add_page_number(slide, page_num)
    return slide


def build_agenda_slide(prs: Presentation, content: SlideContent, page_num: int = 3):
    """Build an agenda slide with highlighted current section."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    items = content.agenda_items or []
    current = content.current_section

    start_top = LAYOUT.CONTENT_TOP + Inches(0.3)
    item_height = Inches(0.7)

    for i, item in enumerate(items):
        top = start_top + i * item_height
        is_current = (current is not None and i == current)

        # Background highlight for current section
        if is_current:
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                LAYOUT.CONTENT_LEFT - Inches(0.1), top - Pt(4),
                LAYOUT.CONTENT_WIDTH + Inches(0.2), item_height - Inches(0.05)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLORS.BACKGROUND_GRAY
            bg.line.fill.background()

        # Number
        txNum = slide.shapes.add_textbox(
            LAYOUT.CONTENT_LEFT, top, Inches(0.5), Inches(0.5)
        )
        _set_text(txNum.text_frame, f"{i+1}.",
                  font_size=Pt(18), bold=True,
                  color=COLORS.MCKINSEY_BLUE if is_current else COLORS.MEDIUM_GRAY)

        # Item text
        txItem = slide.shapes.add_textbox(
            LAYOUT.CONTENT_LEFT + Inches(0.6), top, Inches(10), Inches(0.5)
        )
        _set_text(txItem.text_frame, item,
                  font_size=Pt(16), bold=is_current,
                  color=COLORS.DARK_BLUE if is_current else COLORS.MEDIUM_GRAY)

    _add_page_number(slide, page_num)
    return slide


def build_divider_slide(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(0.3), LAYOUT.HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS.MCKINSEY_BLUE
    bar.line.fill.background()

    # Section number
    if content.section_number is not None:
        txNum = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(2), Inches(1)
        )
        _set_text(txNum.text_frame, f"{content.section_number:02d}",
                  font_name=FONTS.HEADING_FONT, font_size=Pt(48),
                  bold=True, color=COLORS.PALE_GRAY)

    # Section title — use config size
    txTitle = slide.shapes.add_textbox(
        LAYOUT.LEFT_MARGIN * 2, Inches(3.2), Inches(10), Inches(1.5)
    )
    _set_text(txTitle.text_frame, content.action_title,
              font_name=FONTS.HEADING_FONT, font_size=FONTS.TITLE_SIZE,
              bold=FONTS.TITLE_BOLD, color=FONTS.TITLE_COLOR)

    if content.subtitle:
        txSub = slide.shapes.add_textbox(
            LAYOUT.LEFT_MARGIN * 2, Inches(4.5), Inches(10), Inches(0.8)
        )
        _set_text(txSub.text_frame, content.subtitle,
                  font_size=FONTS.REC_LABEL_SIZE, color=COLORS.MEDIUM_GRAY)

    _add_page_number(slide, page_num)
    return slide


def build_content_text(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a text content slide with action title and bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    if content.subtitle:
        txSub = slide.shapes.add_textbox(
            LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP - Inches(0.1),
            LAYOUT.CONTENT_WIDTH, Inches(0.4)
        )
        _set_text(txSub.text_frame, content.subtitle,
                  font_size=FONTS.SUBTITLE_SIZE, bold=True, color=COLORS.MCKINSEY_BLUE)

    # Bullets
    bullet_top = LAYOUT.CONTENT_TOP + (Inches(0.4) if content.subtitle else Inches(0))
    txBox = slide.shapes.add_textbox(
        LAYOUT.CONTENT_LEFT, bullet_top,
        LAYOUT.CONTENT_WIDTH, LAYOUT.CONTENT_HEIGHT
    )
    txBox.text_frame.word_wrap = True
    if content.bullets:
        _add_bullets(txBox.text_frame, content.bullets)

    if content.source:
        _add_source_line(slide, content.source)
    _add_page_number(slide, page_num)
    return slide


def build_content_chart(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a chart slide with action title and full-width chart.

    Uses native editable PowerPoint charts when possible (bar, line, stacked, pie).
    Falls back to matplotlib images for complex charts (waterfall, 2x2, Harvey balls).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    if content.chart:
        # Try native chart first (editable in PowerPoint)
        if USE_NATIVE_CHARTS and can_render_native(content.chart.chart_type):
            add_native_chart(
                slide, content.chart,
                left=LAYOUT.CONTENT_LEFT, top=LAYOUT.CONTENT_TOP,
                width=LAYOUT.CONTENT_WIDTH, height=LAYOUT.CONTENT_HEIGHT,
            )
        else:
            # Fallback to matplotlib image
            chart_image = render_chart(content.chart)
            slide.shapes.add_picture(
                chart_image,
                LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
                LAYOUT.CONTENT_WIDTH, LAYOUT.CONTENT_HEIGHT
            )

    if content.source or (content.chart and content.chart.source):
        src = content.source or content.chart.source
        _add_source_line(slide, src)
    _add_page_number(slide, page_num)
    return slide


def build_content_hybrid(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a two-column slide: text on left, chart on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    # Left column: text/bullets
    txBox = slide.shapes.add_textbox(
        LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
        LAYOUT.COL_LEFT_WIDTH, LAYOUT.CONTENT_HEIGHT
    )
    txBox.text_frame.word_wrap = True
    if content.bullets:
        _add_bullets(txBox.text_frame, content.bullets)

    # Right column: chart (native if possible, image fallback)
    if content.chart:
        if USE_NATIVE_CHARTS and can_render_native(content.chart.chart_type):
            add_native_chart(
                slide, content.chart,
                left=LAYOUT.COL_RIGHT_LEFT, top=LAYOUT.CONTENT_TOP,
                width=LAYOUT.COL_RIGHT_WIDTH, height=LAYOUT.CONTENT_HEIGHT,
            )
        else:
            chart_image = render_chart(content.chart)
            slide.shapes.add_picture(
                chart_image,
                LAYOUT.COL_RIGHT_LEFT, LAYOUT.CONTENT_TOP,
                LAYOUT.COL_RIGHT_WIDTH, LAYOUT.CONTENT_HEIGHT
            )

    if content.source:
        _add_source_line(slide, content.source)
    _add_page_number(slide, page_num)
    return slide


def build_content_table(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a table slide with formatted data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    if content.table:
        rows = len(content.table.rows) + 1  # +1 for header
        cols = len(content.table.headers)

        table_shape = slide.shapes.add_table(
            rows, cols,
            LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
            LAYOUT.CONTENT_WIDTH, min(Inches(4.5), Inches(0.5 * rows))
        )
        table = table_shape.table

        # Header row
        for j, header in enumerate(content.table.headers):
            cell = table.cell(0, j)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = FONTS.BODY_FONT
                    run.font.size = FONTS.TABLE_HEADER_SIZE
                    run.font.bold = True
                    run.font.color.rgb = COLORS.WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS.DARK_BLUE

        # Data rows
        for i, row in enumerate(content.table.rows):
            is_highlight = i in content.table.highlight_rows
            is_last_row = (i == len(content.table.rows) - 1)
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = val

                # Smart alignment: right-align numbers/percentages/currency
                is_numeric = any(c.isdigit() for c in val) and not val[0].isalpha() if val else False

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.RIGHT if is_numeric else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = FONTS.BODY_FONT
                        run.font.size = FONTS.TABLE_BODY_SIZE
                        run.font.color.rgb = COLORS.DARK_GRAY
                        if is_highlight or is_last_row:
                            run.font.bold = True
                if is_highlight:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS.BACKGROUND_GRAY

        # Alternating row shading for readability
        for i in range(len(content.table.rows)):
            if i % 2 == 1 and i not in content.table.highlight_rows:
                for j in range(cols):
                    cell = table.cell(i + 1, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)  # Very subtle gray

        if content.table.source:
            _add_source_line(slide, content.table.source)

    _add_page_number(slide, page_num)
    return slide


def build_recommendation_slide(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a recommendation slide with prioritized items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    if content.bullets:
        txBox = slide.shapes.add_textbox(
            LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
            LAYOUT.CONTENT_WIDTH, LAYOUT.CONTENT_HEIGHT
        )
        txBox.text_frame.word_wrap = True

        for i, bullet in enumerate(content.bullets):
            if i > 0:
                p = txBox.text_frame.add_paragraph()
            else:
                p = txBox.text_frame.paragraphs[0]

            p.space_after = Pt(12)

            # Priority number — use config sizes
            run_num = p.add_run()
            run_num.text = f"{i+1}. "
            run_num.font.name = FONTS.BODY_FONT
            run_num.font.size = FONTS.REC_NUMBER_SIZE
            run_num.font.bold = True
            run_num.font.color.rgb = COLORS.MCKINSEY_BLUE

            # Recommendation text
            if bullet.bold_prefix:
                run_prefix = p.add_run()
                run_prefix.text = bullet.bold_prefix + ": "
                run_prefix.font.name = FONTS.BODY_FONT
                run_prefix.font.size = FONTS.REC_LABEL_SIZE
                run_prefix.font.bold = True
                run_prefix.font.color.rgb = COLORS.DARK_BLUE

            run_text = p.add_run()
            run_text.text = bullet.text
            run_text.font.name = FONTS.BODY_FONT
            run_text.font.size = FONTS.BODY_SIZE
            run_text.font.color.rgb = COLORS.DARK_GRAY

    if content.source:
        _add_source_line(slide, content.source)
    _add_page_number(slide, page_num)
    return slide


def build_next_steps_slide(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a next steps slide with action items table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    items = content.action_items or []
    if items:
        headers = ["Action", "Owner", "Timeline", "Status"]
        rows = len(items) + 1
        cols = 4

        table_shape = slide.shapes.add_table(
            rows, cols,
            LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
            LAYOUT.CONTENT_WIDTH, min(Inches(4.0), Inches(0.5 * rows))
        )
        table = table_shape.table

        # Column widths
        table.columns[0].width = Inches(6.0)
        table.columns[1].width = Inches(2.5)
        table.columns[2].width = Inches(2.0)
        table.columns[3].width = Inches(1.4)

        # Header
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONTS.BODY_FONT
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = COLORS.WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS.DARK_BLUE

        # Data
        for i, item in enumerate(items):
            values = [
                item.get("action", ""),
                item.get("owner", ""),
                item.get("timeline", ""),
                item.get("status", ""),
            ]
            for j, val in enumerate(values):
                cell = table.cell(i + 1, j)
                cell.text = val
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = FONTS.BODY_FONT
                        run.font.size = Pt(10)
                        run.font.color.rgb = COLORS.DARK_GRAY

    _add_page_number(slide, page_num)
    return slide


def build_framework_slide(prs: Presentation, content: SlideContent, page_num: int = 0):
    """Build a framework slide — 2x2 matrix, process flow, or issue tree.

    Detects framework type from content.chart (matrix_2x2) or bullets structure.
    Falls back to a styled text layout if no framework structure is found.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_action_title(slide, content.action_title)

    # Check if it's a 2x2 matrix (chart_type = matrix_2x2)
    chart = content.chart
    if chart and hasattr(chart, 'chart_type') and chart.chart_type == "matrix_2x2":
        _build_2x2_matrix(slide, chart)
    elif content.bullets and len(content.bullets) >= 3:
        _build_process_flow(slide, content.bullets)
    else:
        # Fallback: styled text with bullets
        if content.bullets:
            txBox = slide.shapes.add_textbox(
                LAYOUT.CONTENT_LEFT, LAYOUT.CONTENT_TOP,
                LAYOUT.CONTENT_WIDTH, LAYOUT.CONTENT_HEIGHT,
            )
            txBox.text_frame.word_wrap = True
            _add_bullets(txBox.text_frame, content.bullets)

    if content.source:
        _add_source_line(slide, content.source)
    _add_page_number(slide, page_num)
    return slide


def _build_2x2_matrix(slide, chart):
    """Build a 2x2 matrix visualization on a slide."""
    # Matrix area
    left = LAYOUT.CONTENT_LEFT + Inches(1)
    top = LAYOUT.CONTENT_TOP + Inches(0.3)
    cell_w = Inches(4.8)
    cell_h = Inches(2.2)
    gap = Inches(0.15)

    # Axis labels
    x_label = chart.categories[0] if chart.categories and len(chart.categories) > 0 else "Low"
    x_label_high = chart.categories[1] if chart.categories and len(chart.categories) > 1 else "High"
    y_label = chart.series[0].name if chart.series else "Low"
    y_label_high = chart.series[1].name if chart.series and len(chart.series) > 1 else "High"

    # Quadrant colors and content
    quadrant_colors = [
        COLORS.BACKGROUND_GRAY,       # Top-left (low-low)
        RGBColor(0xEB, 0xF5, 0xFB),   # Top-right (high-low) — light blue
        RGBColor(0xEA, 0xFA, 0xF1),   # Bottom-left (low-high) — light green
        RGBColor(0x00, 0x29, 0x60),    # Bottom-right (high-high) — navy
    ]

    quadrant_text_colors = [
        COLORS.DARK_GRAY,
        COLORS.DARK_BLUE,
        RGBColor(0x27, 0xAE, 0x60),
        COLORS.WHITE,
    ]

    # Get quadrant labels from series data values (if available)
    quadrant_labels = ["", "", "", ""]
    for s in (chart.series or []):
        for i, v in enumerate(s.values[:4]):
            if isinstance(v, str) and i < 4:
                quadrant_labels[i] = v

    # Draw 4 quadrants
    positions = [
        (left, top),                              # Top-left
        (left + cell_w + gap, top),               # Top-right
        (left, top + cell_h + gap),               # Bottom-left
        (left + cell_w + gap, top + cell_h + gap),  # Bottom-right
    ]

    for idx, (x, y) in enumerate(positions):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cell_w, cell_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = quadrant_colors[idx]
        rect.line.color.rgb = COLORS.PALE_GRAY
        rect.line.width = Pt(1)

        # Add label text in center of quadrant
        label = quadrant_labels[idx]
        if label:
            tf = rect.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.name = FONTS.BODY_FONT
            run.font.size = FONTS.REC_LABEL_SIZE
            run.font.bold = True
            run.font.color.rgb = quadrant_text_colors[idx]

    # X-axis label (bottom)
    txX = slide.shapes.add_textbox(
        left + cell_w - Inches(0.5), top + cell_h * 2 + gap + Inches(0.3),
        Inches(3), Inches(0.3),
    )
    _set_text(txX.text_frame, f"{x_label} → {x_label_high}",
              font_size=FONTS.SMALL_SIZE, color=COLORS.MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # Y-axis label (left side)
    txY = slide.shapes.add_textbox(
        left - Inches(1.2), top + cell_h - Inches(0.15),
        Inches(1), Inches(0.3),
    )
    _set_text(txY.text_frame, f"{y_label} → {y_label_high}",
              font_size=FONTS.SMALL_SIZE, color=COLORS.MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # So-what annotation
    if chart.so_what:
        txSoWhat = slide.shapes.add_textbox(
            left, top + cell_h * 2 + gap + Inches(0.7),
            cell_w * 2 + gap, Inches(0.4),
        )
        _set_text(txSoWhat.text_frame, chart.so_what,
                  font_size=FONTS.SMALL_SIZE, color=COLORS.MCKINSEY_BLUE, bold=True)


def _build_process_flow(slide, bullets):
    """Build a horizontal process flow with arrow-connected boxes."""
    n = min(len(bullets), 5)
    total_width = LAYOUT.CONTENT_WIDTH
    box_width = Inches(2.2)
    arrow_width = Inches(0.4)
    available = total_width - (n * box_width + (n - 1) * arrow_width)
    start_x = LAYOUT.CONTENT_LEFT + available / 2
    box_height = Inches(1.8)
    y = LAYOUT.CONTENT_TOP + Inches(1.0)

    for i in range(n):
        x = start_x + i * (box_width + arrow_width)

        # Box
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height,
        )
        rect.fill.solid()
        if i == n - 1:
            rect.fill.fore_color.rgb = COLORS.DARK_BLUE
            text_color = COLORS.WHITE
        else:
            rect.fill.fore_color.rgb = COLORS.BACKGROUND_GRAY
            text_color = COLORS.DARK_GRAY
        rect.line.color.rgb = COLORS.PALE_GRAY

        # Step number
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run_num = p.add_run()
        run_num.text = f"Step {i + 1}\n"
        run_num.font.name = FONTS.BODY_FONT
        run_num.font.size = FONTS.SMALL_SIZE
        run_num.font.bold = True
        run_num.font.color.rgb = COLORS.MCKINSEY_BLUE if i < n - 1 else RGBColor(0x80, 0xBB, 0xE0)

        # Step text
        bullet = bullets[i]
        text = bullet.bold_prefix + " " + bullet.text if bullet.bold_prefix else bullet.text
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = FONTS.BODY_FONT
        run_text.font.size = FONTS.REC_LABEL_SIZE
        run_text.font.bold = False
        run_text.font.color.rgb = text_color

        # Arrow (between boxes, not after last)
        if i < n - 1:
            arrow_x = x + box_width
            arrow_y = y + box_height / 2 - Inches(0.15)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y,
                arrow_width, Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS.MCKINSEY_BLUE
            arrow.line.fill.background()


# =============================================================================
# DISPATCHER
# =============================================================================

SLIDE_BUILDERS = {
    SlideType.TITLE: build_title_slide,
    SlideType.EXECUTIVE_SUMMARY: build_executive_summary,
    SlideType.AGENDA: build_agenda_slide,
    SlideType.DIVIDER: build_divider_slide,
    SlideType.CONTENT_TEXT: build_content_text,
    SlideType.CONTENT_CHART: build_content_chart,
    SlideType.CONTENT_HYBRID: build_content_hybrid,
    SlideType.CONTENT_TABLE: build_content_table,
    SlideType.CONTENT_FRAMEWORK: build_framework_slide,
    SlideType.RECOMMENDATION: build_recommendation_slide,
    SlideType.NEXT_STEPS: build_next_steps_slide,
    SlideType.APPENDIX_DIVIDER: build_divider_slide,
    SlideType.APPENDIX_CONTENT: build_content_text,
}


def build_slide(prs: Presentation, content: SlideContent, page_num: int = 0,
                client: str = "", date: str = "", confidential: bool = True):
    """Build a single slide based on its type."""
    builder = SLIDE_BUILDERS.get(content.slide_type)
    if not builder:
        raise ValueError(f"Unknown slide type: {content.slide_type}")

    if content.slide_type == SlideType.TITLE:
        return builder(prs, content, client=client, date=date, confidential=confidential)
    elif content.slide_type in (SlideType.EXECUTIVE_SUMMARY, SlideType.AGENDA,
                                 SlideType.DIVIDER, SlideType.APPENDIX_DIVIDER):
        return builder(prs, content, page_num=page_num)
    else:
        return builder(prs, content, page_num=page_num)
