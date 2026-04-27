"""
Native PowerPoint chart builder using python-pptx chart objects.
Produces EDITABLE charts (not images) that can be modified in PowerPoint.

Use this instead of chart_builder.py when you need editable charts.
Trade-off: less visual control (no 'so what' annotations, no Harvey balls)
but charts are fully editable in PowerPoint.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .config import COLORS, FONTS, LAYOUT
from .models import ChartSpec, ChartType


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _apply_mckinsey_style(chart):
    """Apply McKinsey visual styling to a native pptx chart."""
    # Font styling
    chart.font.name = FONTS.CHART_FONT
    chart.font.size = Pt(10)
    chart.font.color.rgb = COLORS.DARK_GRAY

    # Legend
    if chart.has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONTS.CHART_FONT

    # Category axis
    if hasattr(chart, 'category_axis'):
        cat_axis = chart.category_axis
        cat_axis.has_major_gridlines = False
        cat_axis.tick_labels.font.size = Pt(9)
        cat_axis.tick_labels.font.name = FONTS.CHART_FONT
        cat_axis.tick_labels.font.color.rgb = COLORS.MEDIUM_GRAY

    # Value axis
    if hasattr(chart, 'value_axis'):
        val_axis = chart.value_axis
        val_axis.has_major_gridlines = True
        val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.name = FONTS.CHART_FONT
        val_axis.tick_labels.font.color.rgb = COLORS.MEDIUM_GRAY


def _apply_series_colors(chart, spec: ChartSpec):
    """Apply McKinsey color palette to chart series."""
    for i, series in enumerate(chart.series):
        if i < len(spec.series) and spec.series[i].color:
            color = _hex_to_rgb(spec.series[i].color)
        else:
            color = COLORS.CHART_SERIES[i % len(COLORS.CHART_SERIES)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color


def add_native_bar_chart(slide, spec: ChartSpec,
                          left=None, top=None, width=None, height=None,
                          stacked: bool = False) -> None:
    """Add an editable bar chart to a slide."""
    left = left or LAYOUT.CONTENT_LEFT
    top = top or LAYOUT.CONTENT_TOP
    width = width or LAYOUT.CONTENT_WIDTH
    height = height or LAYOUT.CONTENT_HEIGHT

    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    if stacked:
        chart_type = XL_CHART_TYPE.COLUMN_STACKED
    elif len(spec.series) > 1:
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    else:
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    chart_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = chart_frame.chart

    _apply_mckinsey_style(chart)
    _apply_series_colors(chart, spec)

    # Data labels
    for series in chart.series:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.font.name = FONTS.CHART_FONT
        series.data_labels.font.color.rgb = COLORS.DARK_GRAY
        series.data_labels.number_format = '#,##0.0'

    chart.has_legend = len(spec.series) > 1


def add_native_bar_horizontal(slide, spec: ChartSpec,
                               left=None, top=None, width=None, height=None) -> None:
    """Add an editable horizontal bar chart."""
    left = left or LAYOUT.CONTENT_LEFT
    top = top or LAYOUT.CONTENT_TOP
    width = width or LAYOUT.CONTENT_WIDTH
    height = height or LAYOUT.CONTENT_HEIGHT

    chart_data = CategoryChartData()
    chart_data.categories = spec.categories
    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _apply_mckinsey_style(chart)
    _apply_series_colors(chart, spec)

    for series in chart.series:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.number_format = '#,##0.0'

    chart.has_legend = len(spec.series) > 1


def add_native_line_chart(slide, spec: ChartSpec,
                           left=None, top=None, width=None, height=None) -> None:
    """Add an editable line chart."""
    left = left or LAYOUT.CONTENT_LEFT
    top = top or LAYOUT.CONTENT_TOP
    width = width or LAYOUT.CONTENT_WIDTH
    height = height or LAYOUT.CONTENT_HEIGHT

    chart_data = CategoryChartData()
    chart_data.categories = spec.categories
    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _apply_mckinsey_style(chart)

    for i, series in enumerate(chart.series):
        color = COLORS.CHART_SERIES[i % len(COLORS.CHART_SERIES)]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.5)
        series.smooth = False
        # Data labels on line chart data points
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.font.name = FONTS.CHART_FONT
        series.data_labels.font.color.rgb = COLORS.DARK_GRAY
        series.data_labels.number_format = '#,##0.0'

    chart.has_legend = len(spec.series) > 1


def add_native_pie_chart(slide, spec: ChartSpec,
                          left=None, top=None, width=None, height=None) -> None:
    """Add an editable pie chart (use sparingly in McKinsey decks)."""
    left = left or LAYOUT.CONTENT_LEFT
    top = top or LAYOUT.CONTENT_TOP
    width = width or LAYOUT.CONTENT_WIDTH
    height = height or LAYOUT.CONTENT_HEIGHT

    chart_data = CategoryChartData()
    chart_data.categories = spec.categories
    if spec.series:
        chart_data.add_series(spec.series[0].name, spec.series[0].values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    _apply_mckinsey_style(chart)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(9)
    data_labels.font.name = FONTS.CHART_FONT
    data_labels.number_format = '0%'
    data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END


# =============================================================================
# DISPATCHER
# =============================================================================

NATIVE_CHART_BUILDERS = {
    ChartType.BAR_VERTICAL: add_native_bar_chart,
    ChartType.GROUPED_BAR: add_native_bar_chart,
    ChartType.STACKED_BAR: lambda slide, spec, **kw: add_native_bar_chart(slide, spec, stacked=True, **kw),
    ChartType.BAR_HORIZONTAL: add_native_bar_horizontal,
    ChartType.LINE: add_native_line_chart,
    ChartType.PIE: add_native_pie_chart,
}

# Chart types that ONLY work as images (no native pptx equivalent)
IMAGE_ONLY_CHARTS = {
    ChartType.WATERFALL,
    ChartType.MATRIX_2X2,
    ChartType.HARVEY_BALLS,
    ChartType.BUBBLE,
    ChartType.SCATTER,
}


def can_render_native(chart_type: ChartType) -> bool:
    """Check if a chart type can be rendered as a native PowerPoint chart."""
    return chart_type in NATIVE_CHART_BUILDERS


def add_native_chart(slide, spec: ChartSpec,
                      left=None, top=None, width=None, height=None) -> bool:
    """Add a native chart if possible. Returns True if successful, False if not supported.

    Guards against empty data that would crash python-pptx.
    """
    builder = NATIVE_CHART_BUILDERS.get(spec.chart_type)
    if not builder:
        return False

    # Guard: ensure we have valid data before creating chart
    if not spec.series or not spec.categories:
        return False
    # Ensure at least one series has values
    has_data = any(s.values for s in spec.series)
    if not has_data:
        return False

    try:
        builder(slide, spec, left=left, top=top, width=width, height=height)
        return True
    except (IndexError, ValueError, TypeError) as e:
        # Fallback gracefully — chart_builder.py will render as image
        print(f"Native chart failed ({spec.chart_type}): {e}")
        return False
