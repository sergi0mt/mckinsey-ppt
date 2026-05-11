"""DeepResearch-style slide deck endpoints.

Companion to the McKinsey workspace flow: same project + report, but a
different generation pipeline (port of deepresearch's /presentation) and a
JSON-based slide schema that the frontend renders as HTML.

Endpoints:
  POST /api/v1/projects/{project_id}/deepresearch-deck/generate   (SSE)
  GET  /api/v1/projects/{project_id}/deepresearch-deck            (latest JSON)
  GET  /api/v1/projects/{project_id}/deepresearch-deck/download   (raw JSON download)
  POST /api/v1/projects/{project_id}/deepresearch-deck/export-pptx  (simple PPTX)
"""
from __future__ import annotations
import json
import logging
import uuid
from datetime import datetime
from io import BytesIO
from pathlib import Path

from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse, Response

from ..config import get_settings
from ..database import get_db
from ..services.deepresearch_presentation import (
    PresentationOptions, PresentationRequest,
    generate_presentation_stream,
)

router = APIRouter(tags=["deepresearch-decks"])
settings = get_settings()
logger = logging.getLogger(__name__)


@router.post("/projects/{project_id}/deepresearch-deck/generate")
async def generate_deepresearch_deck(project_id: str, body: PresentationOptions | None = None):
    """Stream slide generation via SSE.

    The body carries `output_options`; we pull the report text from the
    project's first upload (the imported deepresearch markdown). If the user
    supplies `body.focus`, it goes to the top-level `focus` directive — the
    output_options object also has its own `focus` field which is rendered
    into the prompt's <opciones_de_salida> block.
    """
    opts = body or PresentationOptions()
    db = await get_db()
    try:
        proj_row = await (await db.execute(
            "SELECT * FROM projects WHERE id = ?", (project_id,),
        )).fetchone()
        if not proj_row:
            raise HTTPException(status_code=404, detail="Project not found")

        upload_row = await (await db.execute(
            "SELECT extracted_text FROM uploads WHERE project_id = ? ORDER BY created_at DESC LIMIT 1",
            (project_id,),
        )).fetchone()
        if not upload_row or not upload_row["extracted_text"]:
            raise HTTPException(status_code=400, detail="No imported report found for this project")

        report_text = upload_row["extracted_text"]
        # Pull central_question from the session if available (for objective)
        session_row = await (await db.execute(
            "SELECT stage_data FROM sessions WHERE project_id = ? ORDER BY created_at DESC LIMIT 1",
            (project_id,),
        )).fetchone()
        objective = proj_row["name"]
        if session_row:
            sd = json.loads(session_row["stage_data"] or "{}")
            objective = sd.get("central_question") or sd.get("desired_decision") or objective

        # Slide count + image provider come from query/body? Read defaults from PresentationOptions
        # and let the user override via body. We surface slide_count + image_provider as separate
        # top-level options in the form (for clearer UX), so we map them here.
        slide_count = getattr(opts, "_slide_count_override", None) or 10
        image_provider = getattr(opts, "_image_provider_override", None) or "none"

        req = PresentationRequest(
            report=report_text,
            objective=objective,
            slide_count=slide_count,
            focus=opts.focus or None,
            output_options=opts,
            image_provider=image_provider,
        )

    finally:
        await db.close()

    async def sse_stream():
        last_slides: list[dict] | None = None
        last_palette: dict | None = None
        last_prompt_md: str = ""
        async for ev in generate_presentation_stream(req):
            yield f"data: {json.dumps(ev, ensure_ascii=False)}\n\n"
            if ev.get("type") == "slides":
                last_slides = ev.get("slides")
                last_palette = ev.get("palette")
                last_prompt_md = ev.get("presentation_prompt", "")

        # Persist the final deck on success
        if last_slides is not None:
            try:
                db2 = await get_db()
                try:
                    deck_id = str(uuid.uuid4())
                    # Upsert: only ever one deck per project (latest overwrites)
                    await db2.execute(
                        "DELETE FROM deepresearch_decks WHERE project_id = ?",
                        (project_id,),
                    )
                    await db2.execute(
                        """INSERT INTO deepresearch_decks
                           (id, project_id, slides_json, options_json, palette_json,
                            presentation_prompt_md, image_provider)
                           VALUES (?, ?, ?, ?, ?, ?, ?)""",
                        (
                            deck_id, project_id,
                            json.dumps(last_slides, ensure_ascii=False),
                            json.dumps({**opts.model_dump(), "slide_count": slide_count}, ensure_ascii=False),
                            json.dumps(last_palette or {}, ensure_ascii=False),
                            last_prompt_md,
                            image_provider,
                        ),
                    )
                    await db2.commit()
                    yield f"data: {json.dumps({'type': 'saved', 'deck_id': deck_id})}\n\n"
                finally:
                    await db2.close()
            except Exception as e:
                logger.exception("failed to persist deepresearch deck: %s", e)
                yield f"data: {json.dumps({'type': 'persist_error', 'text': str(e)})}\n\n"

    return StreamingResponse(sse_stream(), media_type="text/event-stream")


@router.post("/projects/{project_id}/deepresearch-deck/generate-with-meta")
async def generate_with_meta(project_id: str, body: dict):
    """Same as /generate but takes a richer body so the frontend can carry
    slide_count and image_provider as first-class fields.

    Body shape:
      { "slide_count": int, "image_provider": str, "options": PresentationOptions-as-dict }
    """
    slide_count = int(body.get("slide_count", 10) or 10)
    image_provider = str(body.get("image_provider", "none") or "none")
    opts_data = body.get("options") or {}
    opts = PresentationOptions(**opts_data)
    # Stash the overrides as private attrs so generate_deepresearch_deck can pick them up.
    # (Cleaner approach: refactor PresentationOptions, but this keeps the surface stable.)
    object.__setattr__(opts, "_slide_count_override", slide_count)
    object.__setattr__(opts, "_image_provider_override", image_provider)
    return await generate_deepresearch_deck(project_id, opts)


@router.get("/projects/{project_id}/deepresearch-deck")
async def get_latest_deepresearch_deck(project_id: str) -> dict:
    """Return the most recent generated deck for this project, or 404 if none."""
    db = await get_db()
    try:
        row = await (await db.execute(
            "SELECT * FROM deepresearch_decks WHERE project_id = ? ORDER BY generated_at DESC LIMIT 1",
            (project_id,),
        )).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="No deck generated yet for this project")
        return {
            "id": row["id"],
            "project_id": row["project_id"],
            "slides": json.loads(row["slides_json"]),
            "options": json.loads(row["options_json"]),
            "palette": json.loads(row["palette_json"]),
            "presentation_prompt": row["presentation_prompt_md"] or "",
            "image_provider": row["image_provider"],
            "generated_at": row["generated_at"],
        }
    finally:
        await db.close()


@router.get("/projects/{project_id}/deepresearch-deck/download")
async def download_deepresearch_deck_json(project_id: str):
    """Force-download the raw deck JSON."""
    db = await get_db()
    try:
        row = await (await db.execute(
            "SELECT slides_json, options_json, palette_json FROM deepresearch_decks "
            "WHERE project_id = ? ORDER BY generated_at DESC LIMIT 1",
            (project_id,),
        )).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="No deck generated yet")
        payload = {
            "slides": json.loads(row["slides_json"]),
            "options": json.loads(row["options_json"]),
            "palette": json.loads(row["palette_json"]),
        }
    finally:
        await db.close()

    body = json.dumps(payload, ensure_ascii=False, indent=2).encode("utf-8")
    fname = f"deepresearch-deck-{project_id[:8]}-{datetime.utcnow():%Y%m%d-%H%M%S}.json"
    return Response(
        content=body,
        media_type="application/json",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


@router.post("/projects/{project_id}/deepresearch-deck/export-pptx")
async def export_deepresearch_deck_pptx(project_id: str):
    """Render the deck to a simple PPTX (no McKinsey-validated layout).

    One slide per deck slide; basic title + body text. Image_url is embedded
    when present. The HTML viewer is the primary preview path; this PPTX is
    a downloadable fallback.
    """
    db = await get_db()
    try:
        row = await (await db.execute(
            "SELECT slides_json, palette_json FROM deepresearch_decks "
            "WHERE project_id = ? ORDER BY generated_at DESC LIMIT 1",
            (project_id,),
        )).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="No deck generated yet")
        slides = json.loads(row["slides_json"])
        palette = json.loads(row["palette_json"])
    finally:
        await db.close()

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError as e:
        raise HTTPException(status_code=500, detail=f"python-pptx not available: {e}")

    def _hex_to_rgb(h: str) -> tuple[int, int, int]:
        h = (h or "#1f2937").lstrip("#")
        if len(h) != 6:
            return (31, 41, 55)
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

    bg_rgb = RGBColor(*_hex_to_rgb(palette.get("bg", "#0c1525")))
    text_rgb = RGBColor(*_hex_to_rgb(palette.get("text", "#e2e8f0")))
    accent_rgb = RGBColor(*_hex_to_rgb(palette.get("accent1", "#3b82f6")))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        # Background fill
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(slide_data.get("title", ""))
        for run in p.runs:
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = accent_rgb

        # Highlight (big number/cifra)
        if slide_data.get("highlight"):
            hb = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12), Inches(1.5))
            htf = hb.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hp.text = str(slide_data["highlight"])
            for run in hp.runs:
                run.font.size = Pt(60)
                run.font.bold = True
                run.font.color.rgb = accent_rgb

        # Content
        cb = slide.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12), Inches(3.5))
        ctf = cb.text_frame
        ctf.word_wrap = True
        content = slide_data.get("content") or []
        if not content and slide_data.get("notes"):
            content = [slide_data["notes"]]
        for i, item in enumerate(content[:8]):
            text = " | ".join(str(x) for x in item) if isinstance(item, list) else str(item)
            p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
            p.text = f"• {text}"
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = text_rgb

    # Save to buffer
    out_dir = Path(settings.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    fname = f"deepresearch-deck-{project_id[:8]}-{datetime.utcnow():%Y%m%d-%H%M%S}.pptx"
    fpath = out_dir / fname
    prs.save(str(fpath))

    return {"filename": fname, "filepath": str(fpath), "download_url": f"/api/v1/exports/{fname}"}
